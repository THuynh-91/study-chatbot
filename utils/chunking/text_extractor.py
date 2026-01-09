from pathlib import Path
import PyPDF2
from pptx import Presentation

# This identifies the doc and calls appropriate text extract
def extract_text(file_path: str) -> list[dict]:
    
    # Grabs the suffix (.pdf, .md, .pptx etc) ... basically the identifier
    extension = Path(file_path).suffix.lower() 

    if extension == ".pdf":
        return extract_pdf(file_path)
    elif extension == ".pptx":
        return extract_pptx(file_path)
    elif extension  in [".txt", ".md"]:
        return extract_text_file(file_path)
    else:
        raise ValueError(f"Unsupported: {extension}")
    
# Extracts text from pdfs
def extract_pdf(file_path: str) -> list[dict]:
    pages_data = []
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)

        for page_num, page in enumerate(pdf_reader.pages, start = 1):
            text = page.extract_text()

            pages_data.append({
                "text": text,
                "page": page_num
            })

    return pages_data

# Extracts text from pptx
def extract_pptx(file_path: str) -> list[dict]:
    pages_data = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start= 1 ):
        text_parts = []

        # Iterate through all shapes on the slides
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)

        pages_data.append({
            "text"  : "\n".join(text_parts),
            "slide": slide_num
        })

    return pages_data

# Plain text files... (already text)
def extract_text_file(file_path: str) -> list[dict]:
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()

    return [{
        "text" : text,
        "page" : 1
    }]
